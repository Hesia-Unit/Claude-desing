"""HESIA pitch deck generator — produces FR and EN .pptx files."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path


STEEL_BLUE = RGBColor(0x0B, 0x25, 0x45)
STEEL_BLUE_LIGHT = RGBColor(0x13, 0x32, 0x57)
SOVEREIGN_GOLD = RGBColor(0xD4, 0xA0, 0x17)
GRAPHITE = RGBColor(0x1A, 0x1E, 0x23)
BONE = RGBColor(0xF5, 0xF1, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x4A, 0x55, 0x68)
LIGHT_SLATE = RGBColor(0x8A, 0x95, 0xA8)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri Light"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = SOVEREIGN_GOLD
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=WHITE,
             bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, lines, size=14, color=WHITE,
                  bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
                  bullet=False, line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        if bullet:
            run.text = f"  {line}"
        else:
            run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_gold_bar(slide, x, y, w=Inches(0.05), h=Inches(0.6)):
    return add_rect(slide, x, y, w, h, SOVEREIGN_GOLD)


def add_footer(slide, page, total, lang):
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), STEEL_BLUE)
    add_text(slide, Inches(0.4), Inches(7.27), Inches(4), Inches(0.22),
             "HESIA — Confidential", size=9, color=LIGHT_SLATE,
             font=FONT_BODY)
    label_inv = "Investor Deck" if lang == "en" else "Pitch Investisseurs"
    add_text(slide, Inches(5), Inches(7.27), Inches(3), Inches(0.22),
             f"{label_inv} · 2026", size=9, color=LIGHT_SLATE,
             font=FONT_BODY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(11), Inches(7.27), Inches(2), Inches(0.22),
             f"{page} / {total}", size=9, color=LIGHT_SLATE,
             font=FONT_BODY, align=PP_ALIGN.RIGHT)


def add_logo(slide, x, y, size=Inches(0.5)):
    hex_outer = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, size, size)
    hex_outer.fill.background()
    hex_outer.line.color.rgb = SOVEREIGN_GOLD
    hex_outer.line.width = Pt(1.5)
    hex_outer.shadow.inherit = False
    tb = slide.shapes.add_textbox(x, y, size, size)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "H"
    run.font.name = FONT_HEAD
    run.font.bold = True
    run.font.size = Pt(int(size.pt * 0.55))
    run.font.color.rgb = SOVEREIGN_GOLD


def add_top_brand(slide):
    add_logo(slide, Inches(0.4), Inches(0.3), Inches(0.45))
    add_text(slide, Inches(0.95), Inches(0.32), Inches(2), Inches(0.4),
             "HESIA", size=16, color=WHITE, bold=True, font=FONT_HEAD)


def add_section_title(slide, kicker, title, color_kicker=SOVEREIGN_GOLD,
                      color_title=WHITE):
    add_text(slide, Inches(0.6), Inches(0.95), Inches(8), Inches(0.35),
             kicker, size=12, color=color_kicker, bold=True,
             font=FONT_HEAD)
    add_text(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
             title, size=32, color=color_title, bold=True,
             font=FONT_HEAD)
    add_rect(slide, Inches(0.6), Inches(2.15), Inches(0.7),
             Inches(0.05), SOVEREIGN_GOLD)


def base_slide(prs, dark=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_color = STEEL_BLUE if dark else BONE
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg_color)
    return slide


def card(slide, x, y, w, h, title, body_lines, accent=SOVEREIGN_GOLD,
         body_size=11):
    add_rect(slide, x, y, w, h, STEEL_BLUE_LIGHT)
    add_rect(slide, x, y, Inches(0.06), h, accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.2),
             w - Inches(0.5), Inches(0.4),
             title, size=14, color=WHITE, bold=True, font=FONT_HEAD)
    add_multiline(slide, x + Inches(0.25), y + Inches(0.7),
                  w - Inches(0.5), h - Inches(0.85),
                  body_lines, size=body_size, color=BONE,
                  font=FONT_BODY, line_spacing=1.35)


def stat_block(slide, x, y, w, h, big, label):
    add_rect(slide, x, y, w, h, STEEL_BLUE_LIGHT)
    add_rect(slide, x, y + h - Inches(0.05), w, Inches(0.05),
             SOVEREIGN_GOLD)
    add_text(slide, x, y + Inches(0.45), w, Inches(0.9),
             big, size=36, color=SOVEREIGN_GOLD, bold=True,
             font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.5), w, Inches(0.6),
             label, size=11, color=BONE,
             font=FONT_BODY, align=PP_ALIGN.CENTER, line_spacing=1.2)


def slide_cover(prs, content):
    slide = base_slide(prs)
    add_rect(slide, Inches(8), 0, Inches(5.333), SLIDE_H, STEEL_BLUE_LIGHT)
    add_rect(slide, Inches(8), Inches(7.45), Inches(5.333), Inches(0.05),
             SOVEREIGN_GOLD)
    add_logo(slide, Inches(0.7), Inches(0.7), Inches(0.6))
    add_text(slide, Inches(1.4), Inches(0.78), Inches(3), Inches(0.5),
             "HESIA", size=22, color=WHITE, bold=True, font=FONT_HEAD)
    add_rect(slide, Inches(0.7), Inches(2.6), Inches(0.7),
             Inches(0.06), SOVEREIGN_GOLD)
    add_text(slide, Inches(0.7), Inches(2.85), Inches(7), Inches(0.6),
             content["cover_kicker"], size=14, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.7), Inches(3.3), Inches(7.3), Inches(2.4),
             content["cover_title"], size=42, color=WHITE, bold=True,
             font=FONT_HEAD, line_spacing=1.05)
    add_text(slide, Inches(0.7), Inches(5.7), Inches(7), Inches(0.5),
             content["cover_subtitle"], size=15, color=BONE,
             font=FONT_BODY)
    add_text(slide, Inches(0.7), Inches(6.6), Inches(7), Inches(0.4),
             content["cover_meta"], size=11, color=LIGHT_SLATE,
             font=FONT_BODY)
    add_text(slide, Inches(8.5), Inches(2.85), Inches(4.5), Inches(0.5),
             content["cover_label"], size=12, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    add_multiline(slide, Inches(8.5), Inches(3.4), Inches(4.5), Inches(3.5),
                  content["cover_highlights"], size=14, color=BONE,
                  font=FONT_BODY, line_spacing=1.6)
    add_text(slide, Inches(8.5), Inches(6.6), Inches(4.5), Inches(0.4),
             content["cover_contact"], size=11, color=LIGHT_SLATE,
             font=FONT_BODY)


def slide_vision(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["vision_kicker"], c["vision_title"])
    add_text(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(2.5),
             c["vision_quote"], size=28, color=WHITE, bold=False,
             font=FONT_HEAD, line_spacing=1.25)
    add_text(slide, Inches(0.6), Inches(5.7), Inches(0.4), Inches(0.4),
             "—", size=20, color=SOVEREIGN_GOLD, bold=True,
             font=FONT_HEAD)
    add_text(slide, Inches(0.95), Inches(5.72), Inches(11), Inches(0.4),
             c["vision_signature"], size=14, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    add_footer(slide, page, total, lang)


def slide_problem(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["problem_kicker"], c["problem_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["problem_intro"], size=15, color=BONE, font=FONT_BODY)
    cards = c["problem_cards"]
    cw = Inches(4.05)
    ch = Inches(3.4)
    gap = Inches(0.2)
    base_x = Inches(0.6)
    for i, (title, lines) in enumerate(cards):
        x = base_x + (cw + gap) * i
        card(slide, x, Inches(3.3), cw, ch, title, lines, body_size=11)
    add_footer(slide, page, total, lang)


def slide_why_now(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["why_kicker"], c["why_title"])
    items = c["why_items"]
    y0 = Inches(2.6)
    for i, (year, title, body) in enumerate(items):
        y = y0 + Inches(0.95) * i
        add_rect(slide, Inches(0.6), y, Inches(1.5), Inches(0.85),
                 STEEL_BLUE_LIGHT)
        add_text(slide, Inches(0.6), y + Inches(0.18), Inches(1.5),
                 Inches(0.5), year, size=18, color=SOVEREIGN_GOLD,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(2.3), y + Inches(0.05), Inches(10.5),
                 Inches(0.4), title, size=15, color=WHITE, bold=True,
                 font=FONT_HEAD)
        add_text(slide, Inches(2.3), y + Inches(0.42), Inches(10.5),
                 Inches(0.5), body, size=12, color=BONE, font=FONT_BODY,
                 line_spacing=1.25)
    add_footer(slide, page, total, lang)


def slide_solution(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["solution_kicker"], c["solution_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["solution_intro"], size=15, color=BONE, font=FONT_BODY)
    layers = c["solution_layers"]
    y0 = Inches(3.2)
    for i, (label, body) in enumerate(layers):
        y = y0 + Inches(0.75) * i
        add_rect(slide, Inches(0.6), y, Inches(3.2), Inches(0.65),
                 STEEL_BLUE_LIGHT)
        add_rect(slide, Inches(0.6), y, Inches(0.06), Inches(0.65),
                 SOVEREIGN_GOLD)
        add_text(slide, Inches(0.85), y + Inches(0.18), Inches(2.8),
                 Inches(0.35), label, size=12, color=WHITE, bold=True,
                 font=FONT_HEAD)
        add_text(slide, Inches(4), y + Inches(0.18), Inches(8.7),
                 Inches(0.4), body, size=12, color=BONE, font=FONT_BODY)
    add_footer(slide, page, total, lang)


def slide_product(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["product_kicker"], c["product_title"])
    modules = c["product_modules"]
    cw = Inches(2.95)
    ch = Inches(4.0)
    gap = Inches(0.18)
    base_x = Inches(0.6)
    for i, m in enumerate(modules):
        x = base_x + (cw + gap) * i
        add_rect(slide, x, Inches(2.7), cw, ch, STEEL_BLUE_LIGHT)
        add_rect(slide, x, Inches(2.7), cw, Inches(0.06), SOVEREIGN_GOLD)
        add_text(slide, x + Inches(0.25), Inches(2.85), cw - Inches(0.5),
                 Inches(0.35), m["status"], size=10,
                 color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), Inches(3.2), cw - Inches(0.5),
                 Inches(0.5), m["name"], size=18, color=WHITE,
                 bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), Inches(3.7), cw - Inches(0.5),
                 Inches(0.45), m["tagline"], size=11,
                 color=SOVEREIGN_GOLD, font=FONT_BODY,
                 line_spacing=1.2)
        add_multiline(slide, x + Inches(0.25), Inches(4.4),
                      cw - Inches(0.5), Inches(2.2), m["features"],
                      size=10, color=BONE, font=FONT_BODY,
                      line_spacing=1.35, bullet=True)
    add_footer(slide, page, total, lang)


def slide_tech(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["tech_kicker"], c["tech_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["tech_intro"], size=15, color=BONE, font=FONT_BODY)
    pillars = c["tech_pillars"]
    cw = Inches(6.05)
    ch = Inches(1.85)
    base_x = Inches(0.6)
    base_y = Inches(3.1)
    for i, (title, lines) in enumerate(pillars):
        col = i % 2
        row = i // 2
        x = base_x + (cw + Inches(0.2)) * col
        y = base_y + (ch + Inches(0.2)) * row
        card(slide, x, y, cw, ch, title, lines, body_size=11)
    add_footer(slide, page, total, lang)


def slide_market(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["market_kicker"], c["market_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["market_intro"], size=15, color=BONE, font=FONT_BODY)
    blocks = c["market_blocks"]
    bw = Inches(4.05)
    bh = Inches(2.4)
    gap = Inches(0.18)
    base_x = Inches(0.6)
    base_y = Inches(3.2)
    for i, (label, big, sub) in enumerate(blocks):
        x = base_x + (bw + gap) * i
        add_rect(slide, x, base_y, bw, bh, STEEL_BLUE_LIGHT)
        add_rect(slide, x, base_y, bw, Inches(0.05), SOVEREIGN_GOLD)
        add_text(slide, x + Inches(0.25), base_y + Inches(0.25),
                 bw - Inches(0.5), Inches(0.4), label, size=12,
                 color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), base_y + Inches(0.7),
                 bw - Inches(0.5), Inches(0.9), big, size=30,
                 color=WHITE, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), base_y + Inches(1.55),
                 bw - Inches(0.5), Inches(0.7), sub, size=11,
                 color=BONE, font=FONT_BODY, line_spacing=1.3)
    add_text(slide, Inches(0.6), Inches(5.85), Inches(12), Inches(1),
             c["market_note"], size=10, color=LIGHT_SLATE,
             font=FONT_BODY, line_spacing=1.4)
    add_footer(slide, page, total, lang)


def slide_gtm(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["gtm_kicker"], c["gtm_title"])
    phases = c["gtm_phases"]
    pw = Inches(6.05)
    ph = Inches(4.0)
    gap = Inches(0.2)
    for i, p in enumerate(phases):
        x = Inches(0.6) + (pw + gap) * i
        add_rect(slide, x, Inches(2.7), pw, ph, STEEL_BLUE_LIGHT)
        add_rect(slide, x, Inches(2.7), pw, Inches(0.05),
                 SOVEREIGN_GOLD)
        add_text(slide, x + Inches(0.3), Inches(2.85), pw - Inches(0.6),
                 Inches(0.35), p["timeline"], size=11,
                 color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.3), Inches(3.2), pw - Inches(0.6),
                 Inches(0.5), p["title"], size=20, color=WHITE,
                 bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.3), Inches(3.75), pw - Inches(0.6),
                 Inches(0.5), p["focus"], size=12,
                 color=BONE, font=FONT_BODY, line_spacing=1.25)
        add_multiline(slide, x + Inches(0.3), Inches(4.6),
                      pw - Inches(0.6), Inches(2.0), p["bullets"],
                      size=11, color=BONE, font=FONT_BODY,
                      line_spacing=1.4, bullet=True)
    add_footer(slide, page, total, lang)


def slide_business(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["biz_kicker"], c["biz_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(7), Inches(0.45),
             c["biz_subtitle1"], size=14, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    rows = c["biz_pricing"]
    table_x = Inches(0.6)
    table_y = Inches(2.85)
    col_widths = [Inches(2.4), Inches(2.0), Inches(2.7)]
    headers = c["biz_headers"]
    add_rect(slide, table_x, table_y,
             sum(col_widths, Emu(0)), Inches(0.4), STEEL_BLUE_LIGHT)
    for i, h in enumerate(headers):
        cx = table_x + sum(col_widths[:i], Emu(0))
        add_text(slide, cx + Inches(0.15), table_y + Inches(0.08),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 h, size=11, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
    for ri, row in enumerate(rows):
        ry = table_y + Inches(0.4) + Inches(0.4) * ri
        if ri % 2 == 0:
            add_rect(slide, table_x, ry,
                     sum(col_widths, Emu(0)), Inches(0.4),
                     STEEL_BLUE_LIGHT)
        for i, val in enumerate(row):
            cx = table_x + sum(col_widths[:i], Emu(0))
            color = WHITE if i == 0 else BONE
            bold = i == 0
            add_text(slide, cx + Inches(0.15), ry + Inches(0.1),
                     col_widths[i] - Inches(0.2), Inches(0.3),
                     val, size=11, color=color, bold=bold,
                     font=FONT_BODY)

    add_text(slide, Inches(8), Inches(2.4), Inches(5), Inches(0.45),
             c["biz_subtitle2"], size=14, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    arr_x = Inches(8)
    arr_y = Inches(2.85)
    arr_w = Inches(4.7)
    arr_h = Inches(3.6)
    add_rect(slide, arr_x, arr_y, arr_w, arr_h, STEEL_BLUE_LIGHT)
    items = c["biz_arr"]
    for i, (label, value) in enumerate(items):
        ry = arr_y + Inches(0.3) + Inches(0.65) * i
        add_text(slide, arr_x + Inches(0.3), ry, Inches(2.3),
                 Inches(0.4), label, size=12, color=BONE,
                 font=FONT_BODY)
        add_text(slide, arr_x + Inches(2.6), ry, Inches(2.0),
                 Inches(0.4), value, size=15, color=SOVEREIGN_GOLD,
                 bold=True, font=FONT_HEAD, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.6), Inches(6.5), Inches(12.2),
             Inches(0.5), c["biz_note"], size=10,
             color=LIGHT_SLATE, font=FONT_BODY)
    add_footer(slide, page, total, lang)


def slide_traction(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["trac_kicker"], c["trac_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["trac_intro"], size=14, color=BONE, font=FONT_BODY)
    quarters = c["trac_quarters"]
    qw = Inches(2.95)
    qh = Inches(3.7)
    gap = Inches(0.18)
    base_x = Inches(0.6)
    base_y = Inches(3.0)
    for i, q in enumerate(quarters):
        x = base_x + (qw + gap) * i
        add_rect(slide, x, base_y, qw, qh, STEEL_BLUE_LIGHT)
        add_rect(slide, x, base_y, qw, Inches(0.05), SOVEREIGN_GOLD)
        add_text(slide, x + Inches(0.2), base_y + Inches(0.2),
                 qw - Inches(0.4), Inches(0.4),
                 q["q"], size=14, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
        add_text(slide, x + Inches(0.2), base_y + Inches(0.6),
                 qw - Inches(0.4), Inches(0.5),
                 q["theme"], size=12, color=WHITE, bold=True,
                 font=FONT_HEAD)
        add_multiline(slide, x + Inches(0.2), base_y + Inches(1.2),
                      qw - Inches(0.4), qh - Inches(1.4),
                      q["items"], size=10, color=BONE,
                      font=FONT_BODY, line_spacing=1.45, bullet=True)
    add_footer(slide, page, total, lang)


def slide_competition(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["comp_kicker"], c["comp_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["comp_intro"], size=14, color=BONE, font=FONT_BODY)
    table_x = Inches(0.6)
    table_y = Inches(3.05)
    headers = c["comp_headers"]
    rows = c["comp_rows"]
    col_widths = [Inches(3.0), Inches(1.85), Inches(2.0),
                  Inches(2.3), Inches(2.95)]
    total_w = sum(col_widths, Emu(0))
    add_rect(slide, table_x, table_y, total_w, Inches(0.45),
             STEEL_BLUE_LIGHT)
    for i, h in enumerate(headers):
        cx = table_x + sum(col_widths[:i], Emu(0))
        add_text(slide, cx + Inches(0.15), table_y + Inches(0.1),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 h, size=11, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
    row_h = Inches(0.5)
    for ri, row in enumerate(rows):
        ry = table_y + Inches(0.45) + row_h * ri
        is_us = "HESIA" in row[0]
        bg = SOVEREIGN_GOLD if is_us else (
            STEEL_BLUE_LIGHT if ri % 2 == 0 else STEEL_BLUE)
        add_rect(slide, table_x, ry, total_w, row_h, bg)
        for i, val in enumerate(row):
            cx = table_x + sum(col_widths[:i], Emu(0))
            if is_us:
                color = STEEL_BLUE
                bold = True
            else:
                color = WHITE if i == 0 else BONE
                bold = i == 0
            add_text(slide, cx + Inches(0.15), ry + Inches(0.13),
                     col_widths[i] - Inches(0.2), row_h - Inches(0.2),
                     val, size=11, color=color, bold=bold,
                     font=FONT_BODY)
    add_footer(slide, page, total, lang)


def slide_team(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["team_kicker"], c["team_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["team_intro"], size=14, color=BONE, font=FONT_BODY)
    members = c["team_members"]
    cw = Inches(4.05)
    ch = Inches(2.6)
    gap = Inches(0.18)
    base_x = Inches(0.6)
    base_y = Inches(3.05)
    for i, m in enumerate(members):
        x = base_x + (cw + gap) * i
        add_rect(slide, x, base_y, cw, ch, STEEL_BLUE_LIGHT)
        add_rect(slide, x, base_y, cw, Inches(0.05), SOVEREIGN_GOLD)
        add_text(slide, x + Inches(0.25), base_y + Inches(0.25),
                 cw - Inches(0.5), Inches(0.4), m["role"], size=11,
                 color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), base_y + Inches(0.6),
                 cw - Inches(0.5), Inches(0.5), m["name"], size=18,
                 color=WHITE, bold=True, font=FONT_HEAD)
        add_text(slide, x + Inches(0.25), base_y + Inches(1.15),
                 cw - Inches(0.5), Inches(1.4), m["bio"], size=11,
                 color=BONE, font=FONT_BODY, line_spacing=1.4)
    advisors_y = Inches(5.95)
    add_text(slide, Inches(0.6), advisors_y, Inches(12), Inches(0.4),
             c["team_advisors_label"], size=12, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.6), advisors_y + Inches(0.4),
             Inches(12), Inches(0.5),
             c["team_advisors"], size=11, color=BONE,
             font=FONT_BODY, line_spacing=1.4)
    add_footer(slide, page, total, lang)


def slide_financials(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["fin_kicker"], c["fin_title"])
    stats = c["fin_stats"]
    sw = Inches(3.0)
    sh = Inches(2.0)
    gap = Inches(0.15)
    base_x = Inches(0.6)
    base_y = Inches(2.5)
    for i, (big, label) in enumerate(stats):
        x = base_x + (sw + gap) * i
        stat_block(slide, x, base_y, sw, sh, big, label)

    add_text(slide, Inches(0.6), Inches(4.85), Inches(7), Inches(0.4),
             c["fin_subtitle"], size=14, color=SOVEREIGN_GOLD,
             bold=True, font=FONT_HEAD)
    rows = c["fin_table"]
    table_y = Inches(5.3)
    col_widths = [Inches(3.0), Inches(2.0), Inches(2.0),
                  Inches(2.0), Inches(3.0)]
    total_w = sum(col_widths, Emu(0))
    add_rect(slide, Inches(0.6), table_y, total_w, Inches(0.4),
             STEEL_BLUE_LIGHT)
    headers = c["fin_table_headers"]
    for i, h in enumerate(headers):
        cx = Inches(0.6) + sum(col_widths[:i], Emu(0))
        add_text(slide, cx + Inches(0.15), table_y + Inches(0.08),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 h, size=11, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
    for ri, row in enumerate(rows):
        ry = table_y + Inches(0.4) + Inches(0.35) * ri
        if ri % 2 == 0:
            add_rect(slide, Inches(0.6), ry, total_w, Inches(0.35),
                     STEEL_BLUE_LIGHT)
        for i, val in enumerate(row):
            cx = Inches(0.6) + sum(col_widths[:i], Emu(0))
            color = WHITE if i == 0 else BONE
            bold = i == 0
            add_text(slide, cx + Inches(0.15), ry + Inches(0.07),
                     col_widths[i] - Inches(0.2), Inches(0.3),
                     val, size=11, color=color, bold=bold,
                     font=FONT_BODY)
    add_footer(slide, page, total, lang)


def slide_use_of_funds(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["uof_kicker"], c["uof_title"])
    add_text(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(0.5),
             c["uof_intro"], size=14, color=BONE, font=FONT_BODY)
    items = c["uof_items"]
    bar_x = Inches(0.6)
    bar_y = Inches(3.2)
    bar_w = Inches(12.1)
    bar_h = Inches(0.5)
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, STEEL_BLUE_LIGHT)
    cumul = 0
    palette = [SOVEREIGN_GOLD, RGBColor(0xC9, 0x90, 0x10),
               RGBColor(0xA0, 0x73, 0x0C), RGBColor(0x70, 0x52, 0x0A),
               RGBColor(0x4A, 0x37, 0x07)]
    for i, (label, pct, _) in enumerate(items):
        seg_w = Inches(bar_w.inches * pct / 100)
        seg_x = bar_x + Inches(bar_w.inches * cumul / 100)
        add_rect(slide, seg_x, bar_y, seg_w, bar_h, palette[i % len(palette)])
        cumul += pct
    y0 = Inches(4.1)
    for i, (label, pct, body) in enumerate(items):
        y = y0 + Inches(0.55) * i
        add_rect(slide, Inches(0.6), y + Inches(0.1), Inches(0.18),
                 Inches(0.18), palette[i % len(palette)])
        add_text(slide, Inches(0.95), y, Inches(3.0), Inches(0.4),
                 label, size=13, color=WHITE, bold=True,
                 font=FONT_HEAD)
        add_text(slide, Inches(4.1), y, Inches(1.2), Inches(0.4),
                 f"{pct}%", size=13, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
        add_text(slide, Inches(5.5), y, Inches(7.5), Inches(0.4),
                 body, size=11, color=BONE, font=FONT_BODY)
    add_footer(slide, page, total, lang)


def slide_ask(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_top_brand(slide)
    add_section_title(slide, c["ask_kicker"], c["ask_title"])
    add_rect(slide, Inches(0.6), Inches(2.6), Inches(8.0),
             Inches(4.0), STEEL_BLUE_LIGHT)
    add_rect(slide, Inches(0.6), Inches(2.6), Inches(0.08),
             Inches(4.0), SOVEREIGN_GOLD)
    add_text(slide, Inches(0.9), Inches(2.85), Inches(7),
             Inches(0.4), c["ask_label"], size=12,
             color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.9), Inches(3.3), Inches(7),
             Inches(1.2), c["ask_amount"], size=64, color=WHITE,
             bold=True, font=FONT_HEAD)
    add_text(slide, Inches(0.9), Inches(4.6), Inches(7),
             Inches(0.5), c["ask_round"], size=18, color=BONE,
             font=FONT_BODY)
    add_multiline(slide, Inches(0.9), Inches(5.3), Inches(7),
                  Inches(1.3), c["ask_terms"], size=11, color=BONE,
                  font=FONT_BODY, line_spacing=1.6, bullet=True)
    add_text(slide, Inches(9.0), Inches(2.85), Inches(4), Inches(0.4),
             c["ask_milestones_label"], size=12,
             color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD)
    items = c["ask_milestones"]
    base_y = Inches(3.3)
    for i, item in enumerate(items):
        y = base_y + Inches(0.55) * i
        add_text(slide, Inches(9.0), y, Inches(0.5), Inches(0.4),
                 "▸", size=14, color=SOVEREIGN_GOLD, bold=True,
                 font=FONT_HEAD)
        add_text(slide, Inches(9.35), y, Inches(3.7), Inches(0.5),
                 item, size=12, color=BONE, font=FONT_BODY,
                 line_spacing=1.3)
    add_footer(slide, page, total, lang)


def slide_closing(prs, c, page, total, lang):
    slide = base_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(2.5), STEEL_BLUE_LIGHT)
    add_logo(slide, Inches(6.4), Inches(2.85), Inches(0.6))
    add_text(slide, Inches(7.05), Inches(2.93), Inches(2),
             Inches(0.5), "HESIA", size=22, color=WHITE, bold=True,
             font=FONT_HEAD)
    add_rect(slide, Inches(6.3), Inches(3.7), Inches(0.7),
             Inches(0.06), SOVEREIGN_GOLD)
    add_text(slide, Inches(0.6), Inches(4.0), Inches(12.1),
             Inches(1.2), c["closing_title"], size=34, color=WHITE,
             bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER,
             line_spacing=1.2)
    add_text(slide, Inches(0.6), Inches(5.4), Inches(12.1),
             Inches(0.6), c["closing_subtitle"], size=15, color=BONE,
             font=FONT_BODY, align=PP_ALIGN.CENTER, line_spacing=1.4)
    add_text(slide, Inches(0.6), Inches(6.4), Inches(12.1),
             Inches(0.4), c["closing_contact"], size=14,
             color=SOVEREIGN_GOLD, bold=True, font=FONT_HEAD,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(6.85), Inches(12.1),
             Inches(0.3), c["closing_meta"], size=10,
             color=LIGHT_SLATE, font=FONT_BODY,
             align=PP_ALIGN.CENTER)


def build_deck(content, out_path, lang):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    builders = [
        ("cover", slide_cover),
        ("vision", slide_vision),
        ("problem", slide_problem),
        ("why_now", slide_why_now),
        ("solution", slide_solution),
        ("product", slide_product),
        ("tech", slide_tech),
        ("market", slide_market),
        ("gtm", slide_gtm),
        ("business", slide_business),
        ("traction", slide_traction),
        ("competition", slide_competition),
        ("team", slide_team),
        ("financials", slide_financials),
        ("use_of_funds", slide_use_of_funds),
        ("ask", slide_ask),
        ("closing", slide_closing),
    ]
    total = len(builders)
    for i, (name, fn) in enumerate(builders, start=1):
        if name == "cover":
            fn(prs, content)
        elif name == "closing":
            fn(prs, content, i, total, lang)
        else:
            fn(prs, content, i, total, lang)
    prs.save(out_path)


def content_fr():
    return {
        "cover_kicker": "PITCH INVESTISSEURS — 2026",
        "cover_title": "Sécurité souveraine pour les systèmes\nautonomes critiques.",
        "cover_subtitle": "Cryptographie post-quantique standardisée. Hardware-backed. Made in Europe.",
        "cover_meta": "HESIA · Paris, France · hesia.eu · Confidentiel",
        "cover_label": "EN UNE PHRASE",
        "cover_highlights": [
            "→ Plateforme firmware + serveur",
            "→ ML-KEM-1024 / ML-DSA-87 / OP-TEE",
            "→ Conforme NIS2 et CRA",
            "→ TRL 4-5 · POC clients en cours",
            "→ Cible : OIV, défense, intégrateurs UE",
        ],
        "cover_contact": "contact@hesia.eu · +33 X XX XX XX XX",

        "vision_kicker": "VISION",
        "vision_title": "Donner à l'Europe le contrôle de ses systèmes autonomes.",
        "vision_quote": "L'Europe ne peut pas piloter sa souveraineté avec des firmwares qu'elle n'a pas écrits, des protocoles qu'elle ne peut pas auditer, et une cryptographie qui sera cassée d'ici 10 ans. HESIA construit la couche logicielle qui rend les drones et systèmes embarqués européens dignes de confiance — aujourd'hui contre les menaces classiques, demain contre l'ordinateur quantique.",
        "vision_signature": "L'équipe fondatrice HESIA",

        "problem_kicker": "LE PROBLÈME",
        "problem_title": "Trois forces convergent — la fenêtre se ferme.",
        "problem_intro": "Les opérateurs européens d'infrastructures critiques et de défense n'ont aujourd'hui aucune solution prête à l'emploi qui combine les trois exigences ci-dessous.",
        "problem_cards": [
            ("Quantum is coming",
             ["NIST a publié les standards FIPS 203/204/205 en août 2024.",
              "L'ANSSI, le BSI, la NSA imposent des roadmaps de migration.",
              "« Harvest now, decrypt later » est un risque actuel, pas futur.",
              "Les drones embarqués sont en retard sur la migration."]),
            ("Réglementation européenne",
             ["NIS2 transposé : audit obligatoire des composants critiques.",
              "Cyber Resilience Act : sécurité par conception, patching, audit.",
              "Sanctions : jusqu'à 2% du CA mondial.",
              "Pas d'outillage fabricant aujourd'hui."]),
            ("Souveraineté & supply chain",
             ["Trop de stacks reposent sur des composants non auditables.",
              "CLOUD Act, FISA 702, ITAR : risque juridique étranger.",
              "Demande EU pour produits indépendants en hausse.",
              "Mais peu d'offre alternative crédible et certifiée."]),
        ],

        "why_kicker": "POURQUOI MAINTENANT",
        "why_title": "Le calendrier réglementaire et technologique aligne tout.",
        "why_items": [
            ("2024",
             "Standards NIST PQC publiés (FIPS 203/204/205)",
             "Premiers standards post-quantiques officiels — la migration peut démarrer industriellement."),
            ("2025",
             "ANSSI guide PQC v2 + BSI directives",
             "Les agences nationales européennes formalisent leurs exigences pour les systèmes critiques."),
            ("2026",
             "NIS2 transposé · CRA s'applique progressivement",
             "Tous les fabricants d'équipements connectés doivent prouver la sécurité par conception."),
            ("2027",
             "CRA pleine application · CSPN/CC EAL4+",
             "Première vague de certifications obligatoires — fenêtre commerciale ouverte."),
        ],

        "solution_kicker": "LA SOLUTION",
        "solution_title": "Une pile complète : du silicium au tableau de bord opérateur.",
        "solution_intro": "HESIA est la première plateforme européenne intégrant cryptographie post-quantique, attestation matérielle et conformité réglementaire pour les drones et systèmes embarqués critiques.",
        "solution_layers": [
            ("Plan opérateur",
             "HESIA Command — supervision flotte, signature policy Ed25519, audit log immuable."),
            ("Plan données / IA",
             "HESIA Observe — pipeline IA embarquée auditable (YOLO + MiDaS), allowlist signée."),
            ("Plan firmware",
             "HESIA Core — TLS 1.3 mTLS, ML-KEM-1024 hybride, sandbox Linux durci."),
            ("Plan matériel",
             "HESIA Attest — TA OP-TEE, scellement clés, attestation distante, secure boot vérifié."),
            ("Conformité",
             "Matrices NIS2 / CRA / ANSSI générées automatiquement, exportables au SIEM."),
        ],

        "product_kicker": "PRODUIT",
        "product_title": "Quatre modules — déployables ensemble ou séparément.",
        "product_modules": [
            {"name": "Core", "status": "DISPONIBLE",
             "tagline": "Firmware durci embarqué.",
             "features": [
                 "TLS 1.3 mTLS + ML-KEM-1024 hybride",
                 "AES-256-GCM rotation < 60s",
                 "Allowlist Ed25519 modules",
                 "Sandbox seccomp / AppArmor",
                 "Audit log chaîné immuable",
                 "Linux Yocto (Jetson, iMX9)"]},
            {"name": "Command", "status": "Q3 2026 ALPHA",
             "tagline": "Supervision multi-drones.",
             "features": [
                 "Console temps réel WebGL",
                 "Signature policy ML-DSA-87",
                 "RBAC + journal opérations",
                 "Connecteur SIEM (Splunk, Elastic)",
                 "API REST / gRPC documentée",
                 "Auto-hébergeable (souverain)"]},
            {"name": "Observe", "status": "Q4 2026 BETA",
             "tagline": "IA embarquée auditable.",
             "features": [
                 "YOLO + MiDaS sandbox",
                 "Limites mémoire / CPU strictes",
                 "Signature des sorties IA",
                 "Modèles vérifiables (hash + sig)",
                 "Profilage comportemental",
                 "Logs RGPD-friendly"]},
            {"name": "Attest", "status": "Q1 2027",
             "tagline": "Attestation matérielle.",
             "features": [
                 "TA OP-TEE / TrustZone",
                 "Scellement clés sur SoC",
                 "Attestation à distance signée",
                 "Secure boot vérifié",
                 "Anti-rollback (RPMB)",
                 "Compatibilité TPM optionnelle"]},
        ],

        "tech_kicker": "DIFFÉRENCIATION TECHNIQUE",
        "tech_title": "Crédibilité technique sourcée — pas de claim sans référence.",
        "tech_intro": "Les choix techniques HESIA reposent sur des standards ouverts, audités et reconnus par les agences nationales.",
        "tech_pillars": [
            ("Cryptographie post-quantique",
             ["• ML-KEM-1024 (NIST FIPS 203, août 2024)",
              "• ML-DSA-87 pour signatures (FIPS 204)",
              "• Mode hybride X25519 + ML-KEM (défense en profondeur)",
              "• Bibliothèque liboqs (Open Quantum Safe — auditée)"]),
            ("Trusted Execution Environment",
             ["• OP-TEE 4.x (Linaro, opensource)",
              "• ARM TrustZone sur Cortex-A",
              "• Scellement matériel des clés maître",
              "• Attestation distante avec exporter binding RFC 9266"]),
            ("Conformité & certification",
             ["• Dossier CSPN ANSSI en instruction",
              "• Gap analysis FIPS 140-3 réalisée",
              "• Roadmap Common Criteria EAL4+ (2027)",
              "• Matrice NIS2 / CRA documentée"]),
            ("Souveraineté technique",
             ["• Code source ouvert aux audits clients",
              "• Hébergement et support 100% UE",
              "• Stack opensource auditable (OP-TEE, OpenSSL, liboqs)",
              "• Pas de dépendance ITAR ni CLOUD Act"]),
        ],

        "market_kicker": "MARCHÉ",
        "market_title": "Un marché en expansion, structurellement souverain.",
        "market_intro": "Sources : Teal Group, PwC, Commission européenne, NIST, ANSSI. Estimations modélisées par HESIA à partir de données publiques.",
        "market_blocks": [
            ("TAM (Total)", "21 Md$",
             "Drones civils + défense + systèmes embarqués critiques · 2030"),
            ("SAM (UE)", "4,2 Md$",
             "Drones d'inspection critique + UAS défense souveraine UE · 2028"),
            ("SOM (5 ans)", "5–25 M€ ARR",
             "Pénétration cible 0,3% du SAM avec offre dual-use civil → défense"),
        ],
        "market_note": "Note : les chiffres ARR sont une cible commerciale prudente. Le SAM est dérivé de la part européenne du marché global drone d'inspection (Teal Group 2024) plus les programmes UAS souverains EU (EDF / OCCAR / DGA budgets publics).",

        "gtm_kicker": "GO-TO-MARKET",
        "gtm_title": "Civil first, défense ensuite — séquence assumée.",
        "gtm_phases": [
            {"timeline": "PHASE 1 · T+0 à T+18 MOIS",
             "title": "Dual-use civil",
             "focus": "Inspection industrielle + infrastructures critiques",
             "bullets": [
                 "Cibles : opérateurs énergie, transport, télécoms",
                 "Cycle : 6-12 mois, ticket 50-200k€",
                 "Effet : track record + crédibilité technique",
                 "Levier : NIS2 / CRA force la décision",
                 "Output : 5-10 clients pilotes UE"]},
            {"timeline": "PHASE 2 · T+18 à T+36 MOIS",
             "title": "Défense souveraine",
             "focus": "Forces armées + primes + intégrateurs",
             "bullets": [
                 "Cibles : DGA, OCCAR, primes EU (Thales, Airbus DS)",
                 "Cycle : 12-24 mois, ticket 200-500k€",
                 "Effet : marges supérieures, contrats pluriannuels",
                 "Levier : track record civil + certifications acquises",
                 "Output : 3-6 contrats défense UE"]},
        ],

        "biz_kicker": "BUSINESS MODEL",
        "biz_title": "Licence SaaS embarquée + services — récurrent à 80%.",
        "biz_subtitle1": "Tarification catalogue",
        "biz_headers": ["Offre", "Tarif", "Métrique"],
        "biz_pricing": [
            ["HESIA Core", "2 500 €", "par drone / an"],
            ["HESIA Command", "15 000 €", "instance / an + sièges"],
            ["HESIA Observe", "1 500 €", "par drone / an"],
            ["HESIA Attest", "800 €", "par drone / an"],
            ["Audit / Assessment", "20 000 €", "forfait projet"],
            ["Support Gold (SLA 24/7)", "20 % MRR", "annuel"],
        ],
        "biz_subtitle2": "Projection ARR",
        "biz_arr": [
            ("Année 1 (2026)", "0,8 – 1,5 M€"),
            ("Année 2 (2027)", "2,5 – 5 M€"),
            ("Année 3 (2028)", "6 – 12 M€"),
            ("Année 5 (2030)", "15 – 30 M€"),
        ],
        "biz_note": "Hypothèses : 5-10 clients année 1 (POC + licence), expansion volume + modules en année 2-3, entrée défense en année 3-4. Conservateur vs leaders catégorie cyber B2B.",

        "trac_kicker": "ROADMAP & TRACTION",
        "trac_title": "Plan d'exécution sur 24 mois — milestones publics.",
        "trac_intro": "Roadmap à exécution disciplinée. Chaque trimestre a un livrable produit + un livrable commercial.",
        "trac_quarters": [
            {"q": "T1 2026", "theme": "Lancement public",
             "items": ["HESIA Core 1.0",
                       "Site web + RP",
                       "Premier livre blanc PQC",
                       "1er client pilote signé",
                       "Présence FIC (Lille)"]},
            {"q": "T2 2026", "theme": "Pre-CSPN",
             "items": ["Core 1.1 + audit externe",
                       "OP-TEE TA alpha",
                       "Dossier CSPN déposé",
                       "2-3 clients pilotes additionnels",
                       "Levée de fonds (objectif)"]},
            {"q": "T3 2026", "theme": "Command beta",
             "items": ["HESIA Command alpha",
                       "Connecteurs SIEM",
                       "Recrutement Security Engineer",
                       "UAV Show (Bordeaux)",
                       "1ère étude de cas publique"]},
            {"q": "T4 2026", "theme": "Maturité produit",
             "items": ["HESIA Attest beta",
                       "Common Criteria gap",
                       "Bilan année + ARR ≥ 800k€",
                       "Milipol Paris (stand)",
                       "Préparation expansion DACH"]},
        ],

        "comp_kicker": "CONCURRENCE",
        "comp_title": "Notre place sur l'échiquier — un trou dans le marché.",
        "comp_intro": "HESIA n'est pas un fabricant de drones, ni un cabinet de conseil cyber. C'est l'éditeur logiciel souverain manquant.",
        "comp_headers": ["Acteur", "Origine", "Approche",
                         "PQC natif", "Souveraineté UE"],
        "comp_rows": [
            ["Skydio Defense", "🇺🇸 USA", "HW + SW intégré",
             "Roadmap", "Non (ITAR)"],
            ["Anduril", "🇺🇸 USA", "HW + SW + IA",
             "Roadmap", "Non"],
            ["Parrot", "🇫🇷 France", "HW grand public + pro",
             "Partiel", "Partielle"],
            ["Quantum Systems", "🇩🇪 Allemagne", "HW reconnaissance",
             "Pas de claim public", "Oui"],
            ["Elbit / IAI", "🇮🇱 Israël", "HW défense intégré",
             "Roadmap", "Non"],
            ["HESIA", "🇫🇷 → 🇪🇺", "SW horizontal embarqué",
             "Natif", "Native"],
        ],

        "team_kicker": "ÉQUIPE",
        "team_title": "Profils techniques rares — combinaison unique cyber + embarqué + UE.",
        "team_intro": "Équipe fondatrice de [N] personnes, complétée par freelances spécialisés et conseil scientifique externe.",
        "team_members": [
            {"role": "CO-FONDATEUR & CTO",
             "name": "[À compléter]",
             "bio": "Expérience cryptographie + sécurité embarquée. Parcours [ANSSI / Thales / Airbus / labo recherche]. Expert OP-TEE et migration PQC. Publications publiques sur [sujets]."},
            {"role": "CO-FONDATEUR & CEO",
             "name": "[À compléter]",
             "bio": "Expérience produit + business B2B deeptech. Parcours [scale-up cyber / défense / industriel]. Réseau acheteurs OIV et défense UE."},
            {"role": "HEAD OF SECURITY ENGINEERING",
             "name": "[À recruter Q3 2026]",
             "bio": "Cible profil senior 8-12 ans : OP-TEE, TrustZone, Yocto, Linux kernel. Expérience cycle certification CSPN ou Common Criteria souhaitée."},
        ],
        "team_advisors_label": "CONSEIL & ADVISORS",
        "team_advisors": "Conseil scientifique (cryptographie / défense) · Advisor go-to-market défense · Advisor compliance NIS2 / CRA · Mentors French Tech / EIC.",

        "fin_kicker": "FINANCIER",
        "fin_title": "Modèle financier — projections conservatrices.",
        "fin_stats": [
            ("0,8–1,5 M€", "ARR cible année 1"),
            ("65–80 %", "Marge brute SaaS embarqué"),
            ("9–12 mois", "Cycle de vente moyen"),
            ("3–5×", "LTV/CAC à maturité"),
        ],
        "fin_subtitle": "Trajectoire 5 ans (scénario médian)",
        "fin_table_headers": ["Indicateur (M€)", "2026", "2027",
                              "2028", "2030 cible"],
        "fin_table": [
            ["ARR signé cumulé", "1,0", "3,5", "9,0", "22,0"],
            ["Clients payants", "5-8", "12-18", "25-35", "60+"],
            ["Effectif", "10", "20", "35", "70"],
            ["Burn / mois", "0,1", "0,2", "0,3", "0,4"],
            ["Marge brute", "60%", "70%", "75%", "80%"],
            ["Cash runway", ">18m", ">18m", ">18m", "auto-financé"],
        ],

        "uof_kicker": "UTILISATION DES FONDS",
        "uof_title": "Allocation des fonds levés — 18-24 mois de runway.",
        "uof_intro": "Priorité : produit certifié + équipe technique + traction commerciale UE. Pas de dépenses brand prématurées.",
        "uof_items": [
            ("R&D produit + certification", 45,
             "Recrutement ingénierie (5-7 ETP), CSPN, audit pentest externe, gap CC EAL4+."),
            ("Sales & marketing", 25,
             "AE senior + SE + content marketing + 2 salons stand UE."),
            ("Opérations & équipe", 15,
             "CFO fractional, juridique, RH, outils, infrastructure."),
            ("Buffer stratégique", 10,
             "Réserve opportunités M&A, partenariats stratégiques."),
            ("Frais juridiques & gouvernance", 5,
             "Levée, pacte associés, propriété intellectuelle, conformité."),
        ],

        "ask_kicker": "ASK",
        "ask_title": "Notre demande — partenariat sur 18 mois.",
        "ask_label": "MONTANT RECHERCHÉ",
        "ask_amount": "5 – 8 M€",
        "ask_round": "Série Seed+ / Série A",
        "ask_terms": [
            "Lead investor recherché : 2-4 M€",
            "Co-lead souhaité : 1-2 M€",
            "Tickets stratégiques : 0,5-1 M€",
            "Dispositifs publics complémentaires : Bpifrance, EIC, France 2030",
            "Préférence investisseurs UE / souverains",
        ],
        "ask_milestones_label": "MILESTONES SUR 18 MOIS",
        "ask_milestones": [
            "ARR signé ≥ 3 M€ d'ici fin 2027",
            "Certification CSPN obtenue",
            "OP-TEE TA en production",
            "10+ clients pilotes UE",
            "1er contrat défense souveraine signé",
            "Préparation Série A (12-25 M€)",
        ],

        "closing_title": "Sécurité souveraine.\nPost-quantique. Aujourd'hui.",
        "closing_subtitle": "Échanger avec l'équipe fondatrice : nous sommes ouverts aux discussions stratégiques, partenariats techniques et investissements.",
        "closing_contact": "investors@hesia.eu  ·  hesia.eu",
        "closing_meta": "Document confidentiel — usage limité aux destinataires identifiés. © HESIA 2026.",
    }


def content_en():
    return {
        "cover_kicker": "INVESTOR DECK — 2026",
        "cover_title": "Sovereign security for critical\nautonomous systems.",
        "cover_subtitle": "Standardized post-quantum cryptography. Hardware-backed. Made in Europe.",
        "cover_meta": "HESIA · Paris, France · hesia.eu · Confidential",
        "cover_label": "AT A GLANCE",
        "cover_highlights": [
            "→ Firmware + server platform",
            "→ ML-KEM-1024 / ML-DSA-87 / OP-TEE",
            "→ NIS2 and CRA compliant",
            "→ TRL 4-5 · Customer pilots in flight",
            "→ Targets: critical operators, EU defense",
        ],
        "cover_contact": "contact@hesia.eu · +33 X XX XX XX XX",

        "vision_kicker": "VISION",
        "vision_title": "Give Europe control over its autonomous systems.",
        "vision_quote": "Europe cannot steer its sovereignty with firmwares it did not write, protocols it cannot audit, and cryptography that will be broken within a decade. HESIA builds the software layer that makes European drones and embedded systems trustworthy — today against classical threats, tomorrow against the quantum computer.",
        "vision_signature": "The HESIA founding team",

        "problem_kicker": "THE PROBLEM",
        "problem_title": "Three forces converge — the window is closing.",
        "problem_intro": "European critical infrastructure operators and defense actors have no off-the-shelf solution today that meets all three requirements below.",
        "problem_cards": [
            ("Quantum is coming",
             ["NIST published FIPS 203/204/205 in August 2024.",
              "ANSSI, BSI, NSA mandate migration roadmaps.",
              "\"Harvest now, decrypt later\" is a present risk.",
              "Embedded drones lag on PQC migration."]),
            ("EU regulation",
             ["NIS2 transposed: mandatory audit of critical components.",
              "Cyber Resilience Act: security-by-design, patching, audits.",
              "Penalties: up to 2% of global revenue.",
              "No off-the-shelf manufacturer toolkit today."]),
            ("Sovereignty & supply chain",
             ["Most stacks rely on non-auditable components.",
              "CLOUD Act, FISA 702, ITAR: foreign legal exposure.",
              "Rising EU demand for independent products.",
              "Few credible certified alternatives."]),
        ],

        "why_kicker": "WHY NOW",
        "why_title": "Regulation and technology align — perfect window.",
        "why_items": [
            ("2024",
             "NIST PQC standards published (FIPS 203/204/205)",
             "First official post-quantum standards — industrial migration can begin."),
            ("2025",
             "ANSSI PQC guide v2 + BSI directives",
             "European national agencies formalize requirements for critical systems."),
            ("2026",
             "NIS2 transposed · CRA enters into force progressively",
             "All connected-device makers must demonstrate security-by-design."),
            ("2027",
             "CRA full enforcement · CSPN/CC EAL4+ wave",
             "First mandatory certifications — commercial window opens wide."),
        ],

        "solution_kicker": "THE SOLUTION",
        "solution_title": "Full stack: from silicon to operator dashboard.",
        "solution_intro": "HESIA is the first European platform integrating post-quantum cryptography, hardware attestation, and regulatory compliance for critical drones and embedded systems.",
        "solution_layers": [
            ("Operator plane",
             "HESIA Command — fleet supervision, Ed25519 policy signing, immutable audit log."),
            ("Data / AI plane",
             "HESIA Observe — auditable embedded AI pipeline (YOLO + MiDaS), signed allowlist."),
            ("Firmware plane",
             "HESIA Core — TLS 1.3 mTLS, hybrid ML-KEM-1024, hardened Linux sandbox."),
            ("Hardware plane",
             "HESIA Attest — OP-TEE TA, key sealing, remote attestation, verified secure boot."),
            ("Compliance",
             "Auto-generated NIS2 / CRA / ANSSI matrices, exportable to SIEM."),
        ],

        "product_kicker": "PRODUCT",
        "product_title": "Four modules — deployable together or independently.",
        "product_modules": [
            {"name": "Core", "status": "AVAILABLE",
             "tagline": "Hardened embedded firmware.",
             "features": [
                 "TLS 1.3 mTLS + hybrid ML-KEM-1024",
                 "AES-256-GCM rotation < 60s",
                 "Ed25519 module allowlist",
                 "seccomp / AppArmor sandbox",
                 "Chained immutable audit log",
                 "Linux Yocto (Jetson, iMX9)"]},
            {"name": "Command", "status": "Q3 2026 ALPHA",
             "tagline": "Multi-drone supervision.",
             "features": [
                 "Real-time WebGL console",
                 "ML-DSA-87 policy signing",
                 "RBAC + operations log",
                 "SIEM connector (Splunk, Elastic)",
                 "Documented REST / gRPC API",
                 "Self-hostable (sovereign)"]},
            {"name": "Observe", "status": "Q4 2026 BETA",
             "tagline": "Auditable embedded AI.",
             "features": [
                 "YOLO + MiDaS sandbox",
                 "Strict memory / CPU limits",
                 "Signed AI outputs",
                 "Verifiable models (hash + sig)",
                 "Behavior profiling",
                 "GDPR-friendly logs"]},
            {"name": "Attest", "status": "Q1 2027",
             "tagline": "Hardware-backed attestation.",
             "features": [
                 "OP-TEE / TrustZone TA",
                 "Key sealing on SoC",
                 "Signed remote attestation",
                 "Verified secure boot",
                 "Anti-rollback (RPMB)",
                 "Optional TPM compatibility"]},
        ],

        "tech_kicker": "TECHNICAL DIFFERENTIATION",
        "tech_title": "Sourced technical credibility — no claim without reference.",
        "tech_intro": "HESIA's technical choices rely on open, audited standards recognized by national agencies.",
        "tech_pillars": [
            ("Post-quantum cryptography",
             ["• ML-KEM-1024 (NIST FIPS 203, August 2024)",
              "• ML-DSA-87 for signatures (FIPS 204)",
              "• Hybrid mode X25519 + ML-KEM (defense in depth)",
              "• liboqs library (Open Quantum Safe — audited)"]),
            ("Trusted Execution Environment",
             ["• OP-TEE 4.x (Linaro, open source)",
              "• ARM TrustZone on Cortex-A",
              "• Hardware-backed master key sealing",
              "• Remote attestation with RFC 9266 exporter binding"]),
            ("Compliance & certification",
             ["• ANSSI CSPN application under review",
              "• FIPS 140-3 gap analysis completed",
              "• Common Criteria EAL4+ roadmap (2027)",
              "• NIS2 / CRA documented mapping"]),
            ("Technical sovereignty",
             ["• Source code open to customer audits",
              "• 100% EU hosting and support",
              "• Auditable open-source stack (OP-TEE, OpenSSL, liboqs)",
              "• No ITAR or CLOUD Act dependency"]),
        ],

        "market_kicker": "MARKET",
        "market_title": "An expanding market — structurally sovereign.",
        "market_intro": "Sources: Teal Group, PwC, European Commission, NIST, ANSSI. Estimates modeled by HESIA from public data.",
        "market_blocks": [
            ("TAM (Total)", "$21B",
             "Civil + defense drones + critical embedded systems · 2030"),
            ("SAM (EU)", "$4.2B",
             "Critical inspection drones + EU sovereign UAS defense · 2028"),
            ("SOM (5y)", "€5–25M ARR",
             "Target 0.3% SAM penetration · dual-use civil → defense"),
        ],
        "market_note": "Note: ARR figures are a conservative commercial target. SAM derives from European share of global inspection drone market (Teal Group 2024) plus EU sovereign UAS programs (EDF / OCCAR / DGA public budgets).",

        "gtm_kicker": "GO-TO-MARKET",
        "gtm_title": "Civil first, defense next — explicit sequence.",
        "gtm_phases": [
            {"timeline": "PHASE 1 · M+0 TO M+18",
             "title": "Civil dual-use",
             "focus": "Industrial inspection + critical infrastructure",
             "bullets": [
                 "Targets: energy, transport, telecom operators",
                 "Cycle: 6-12 months, deal size €50-200k",
                 "Effect: track record + technical credibility",
                 "Driver: NIS2 / CRA forces decision",
                 "Output: 5-10 EU pilot customers"]},
            {"timeline": "PHASE 2 · M+18 TO M+36",
             "title": "Sovereign defense",
             "focus": "Armed forces + primes + integrators",
             "bullets": [
                 "Targets: DGA, OCCAR, EU primes (Thales, Airbus DS)",
                 "Cycle: 12-24 months, deal size €200-500k",
                 "Effect: higher margins, multi-year contracts",
                 "Driver: civil track record + certifications",
                 "Output: 3-6 EU defense contracts"]},
        ],

        "biz_kicker": "BUSINESS MODEL",
        "biz_title": "Embedded SaaS license + services — 80% recurring.",
        "biz_subtitle1": "List pricing",
        "biz_headers": ["Offer", "Price", "Metric"],
        "biz_pricing": [
            ["HESIA Core", "€2,500", "per drone / year"],
            ["HESIA Command", "€15,000", "instance / year + seats"],
            ["HESIA Observe", "€1,500", "per drone / year"],
            ["HESIA Attest", "€800", "per drone / year"],
            ["Audit / Assessment", "€20,000", "fixed project"],
            ["Gold Support (24/7 SLA)", "20% MRR", "annual"],
        ],
        "biz_subtitle2": "ARR projection",
        "biz_arr": [
            ("Year 1 (2026)", "€0.8 – 1.5M"),
            ("Year 2 (2027)", "€2.5 – 5M"),
            ("Year 3 (2028)", "€6 – 12M"),
            ("Year 5 (2030)", "€15 – 30M"),
        ],
        "biz_note": "Assumptions: 5-10 customers Y1 (POC + license), volume + module expansion Y2-Y3, defense entry Y3-Y4. Conservative vs B2B cyber category leaders.",

        "trac_kicker": "ROADMAP & TRACTION",
        "trac_title": "24-month execution plan — public milestones.",
        "trac_intro": "Disciplined roadmap. Each quarter delivers a product milestone + a commercial milestone.",
        "trac_quarters": [
            {"q": "Q1 2026", "theme": "Public launch",
             "items": ["HESIA Core 1.0",
                       "Website + PR",
                       "First PQC whitepaper",
                       "1st pilot customer signed",
                       "FIC presence (Lille)"]},
            {"q": "Q2 2026", "theme": "Pre-CSPN",
             "items": ["Core 1.1 + external audit",
                       "OP-TEE TA alpha",
                       "CSPN application filed",
                       "2-3 additional pilots",
                       "Funding round (target)"]},
            {"q": "Q3 2026", "theme": "Command beta",
             "items": ["HESIA Command alpha",
                       "SIEM connectors",
                       "Hire Security Engineer",
                       "UAV Show (Bordeaux)",
                       "First public case study"]},
            {"q": "Q4 2026", "theme": "Product maturity",
             "items": ["HESIA Attest beta",
                       "Common Criteria gap",
                       "Year review · ARR ≥ €800k",
                       "Milipol Paris (booth)",
                       "DACH expansion prep"]},
        ],

        "comp_kicker": "COMPETITION",
        "comp_title": "Where we sit on the board — a real market gap.",
        "comp_intro": "HESIA is not a drone manufacturer, nor a cyber consultancy. It is the missing sovereign software vendor.",
        "comp_headers": ["Player", "Origin", "Approach",
                         "Native PQC", "EU sovereignty"],
        "comp_rows": [
            ["Skydio Defense", "🇺🇸 USA", "Integrated HW + SW",
             "Roadmap", "No (ITAR)"],
            ["Anduril", "🇺🇸 USA", "HW + SW + AI",
             "Roadmap", "No"],
            ["Parrot", "🇫🇷 France", "Consumer + pro HW",
             "Partial", "Partial"],
            ["Quantum Systems", "🇩🇪 Germany", "Reconnaissance HW",
             "No public claim", "Yes"],
            ["Elbit / IAI", "🇮🇱 Israel", "Integrated defense HW",
             "Roadmap", "No"],
            ["HESIA", "🇫🇷 → 🇪🇺", "Horizontal embedded SW",
             "Native", "Native"],
        ],

        "team_kicker": "TEAM",
        "team_title": "Rare technical profiles — unique cyber + embedded + EU mix.",
        "team_intro": "Founding team of [N] people, complemented by specialized freelancers and external scientific board.",
        "team_members": [
            {"role": "CO-FOUNDER & CTO",
             "name": "[To be filled]",
             "bio": "Cryptography + embedded security background. Previously at [ANSSI / Thales / Airbus / research lab]. OP-TEE expert and PQC migration. Public publications on [topics]."},
            {"role": "CO-FOUNDER & CEO",
             "name": "[To be filled]",
             "bio": "Product + B2B deeptech business background. Previously at [cyber / defense / industrial scale-up]. Network in EU OIV and defense procurement."},
            {"role": "HEAD OF SECURITY ENGINEERING",
             "name": "[To hire Q3 2026]",
             "bio": "Targeting senior 8-12 years: OP-TEE, TrustZone, Yocto, Linux kernel. CSPN or Common Criteria certification cycle experience preferred."},
        ],
        "team_advisors_label": "ADVISORS",
        "team_advisors": "Scientific board (cryptography / defense) · Defense GTM advisor · NIS2 / CRA compliance advisor · French Tech / EIC mentors.",

        "fin_kicker": "FINANCIALS",
        "fin_title": "Financial model — conservative projections.",
        "fin_stats": [
            ("€0.8–1.5M", "Year 1 ARR target"),
            ("65–80%", "Embedded SaaS gross margin"),
            ("9–12 months", "Average sales cycle"),
            ("3–5×", "LTV/CAC at maturity"),
        ],
        "fin_subtitle": "5-year trajectory (median scenario)",
        "fin_table_headers": ["Indicator (€M)", "2026", "2027",
                              "2028", "2030 target"],
        "fin_table": [
            ["Cumulative signed ARR", "1.0", "3.5", "9.0", "22.0"],
            ["Paying customers", "5-8", "12-18", "25-35", "60+"],
            ["Headcount", "10", "20", "35", "70"],
            ["Burn / month", "0.1", "0.2", "0.3", "0.4"],
            ["Gross margin", "60%", "70%", "75%", "80%"],
            ["Cash runway", ">18m", ">18m", ">18m", "self-funded"],
        ],

        "uof_kicker": "USE OF FUNDS",
        "uof_title": "Allocation of raised funds — 18-24 months runway.",
        "uof_intro": "Priority: certified product + technical team + EU commercial traction. No premature brand spending.",
        "uof_items": [
            ("Product R&D + certification", 45,
             "Engineering hires (5-7 FTE), CSPN, external pentest audit, CC EAL4+ gap."),
            ("Sales & marketing", 25,
             "Senior AE + SE + content marketing + 2 EU trade-show booths."),
            ("Operations & team", 15,
             "Fractional CFO, legal, HR, tooling, infrastructure."),
            ("Strategic buffer", 10,
             "Reserve for M&A opportunities, strategic partnerships."),
            ("Legal & governance", 5,
             "Round, shareholders agreement, IP, compliance."),
        ],

        "ask_kicker": "ASK",
        "ask_title": "Our ask — 18-month partnership.",
        "ask_label": "AMOUNT SOUGHT",
        "ask_amount": "€5 – 8M",
        "ask_round": "Seed+ / Series A",
        "ask_terms": [
            "Lead investor: €2-4M",
            "Co-lead: €1-2M",
            "Strategic tickets: €0.5-1M",
            "Complementary public schemes: Bpifrance, EIC, France 2030",
            "Preference for EU / sovereign investors",
        ],
        "ask_milestones_label": "18-MONTH MILESTONES",
        "ask_milestones": [
            "Signed ARR ≥ €3M by end 2027",
            "CSPN certification obtained",
            "OP-TEE TA in production",
            "10+ EU pilot customers",
            "1st sovereign defense contract signed",
            "Series A preparation (€12-25M)",
        ],

        "closing_title": "Sovereign security.\nPost-quantum. Now.",
        "closing_subtitle": "Reach out to the founding team — we welcome strategic discussions, technical partnerships, and investments.",
        "closing_contact": "investors@hesia.eu  ·  hesia.eu",
        "closing_meta": "Confidential — for identified recipients only. © HESIA 2026.",
    }


def main():
    out_dir = Path(__file__).parent
    fr_path = out_dir / "HESIA_Pitch_FR.pptx"
    en_path = out_dir / "HESIA_Pitch_EN.pptx"
    build_deck(content_fr(), fr_path, "fr")
    build_deck(content_en(), en_path, "en")
    print(f"Generated:\n  {fr_path}\n  {en_path}")


if __name__ == "__main__":
    main()
